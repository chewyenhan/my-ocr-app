import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

class GlobalOffsets:
    """
    Global Calibration Class to fine-tune offsets dynamically.
    Widescreen PowerPoint slides are configured as standard 16:9 (10.0 x 5.625 inches).
    """
    def __init__(self, offset_x=0.0, offset_y=0.0, scale_x=1.0, scale_y=1.0):
        self.offset_x = offset_x  # shift horizontal in inches
        self.offset_y = offset_y  # shift vertical in inches
        self.scale_x = scale_x    # width scale coefficient
        self.scale_y = scale_y    # height scale coefficient

    def adjust_x(self, x_percent):
        base_width = 10.0
        val = (float(x_percent) / 100.0) * base_width * self.scale_x + self.offset_x
        return Inches(max(0.0, min(base_width, val)))

    def adjust_y(self, y_percent):
        base_height = 5.625
        val = (float(y_percent) / 100.0) * base_height * self.scale_y + self.offset_y
        return Inches(max(0.0, min(base_height, val)))

    def adjust_width(self, w_percent):
        base_width = 10.0
        val = (float(w_percent) / 100.0) * base_width * self.scale_x
        return Inches(max(0.1, min(base_width, val)))

    def adjust_height(self, h_percent):
        base_height = 5.625
        val = (float(h_percent) / 100.0) * base_height * self.scale_y
        return Inches(max(0.1, min(base_height, val)))


class PPTRebuildEngine:
    """
    Architectural rebuild engine compiling layout layers:
    1. Background Layer (Solid color)
    2. Vector Layer (Visual cards, shapes, dividers, and rounded rect blocks)
    3. Typography Text Layer (Flexible text boxes with custom alignment)
    """
    def __init__(self, offsets: GlobalOffsets):
        self.offsets = offsets

    def _hex_to_rgb(self, hex_str: str) -> RGBColor:
        hex_str = hex_str.strip('#')
        if len(hex_str) != 6:
            return RGBColor(240, 240, 240)  # Default light color fallback
        try:
            r = int(hex_str[0:2], 16)
            g = int(hex_str[2:4], 16)
            b = int(hex_str[4:6], 16)
            return RGBColor(r, g, b)
        except ValueError:
            return RGBColor(240, 240, 240)

    def rebuild_presentation(self, slide_definitions: list) -> Presentation:
        prs = Presentation()
        # Set widescreen layout 16x9 default size explicitly (10 x 5.625 inches)
        prs.slide_width = Inches(10.0)
        prs.slide_height = Inches(5.625)

        # Clear standard layout templates to work on clean pages
        blank_slide_layout = prs.slide_layouts[6] 

        for s_index, slide_def in enumerate(slide_definitions):
            slide = prs.slides.add_slide(blank_slide_layout)

            # LAYER 1: Background Layer Rebuilding
            bg_def = slide_def.get("background", {})
            bg_color_hex = bg_def.get("color", "FFFFFF")
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = self._hex_to_rgb(bg_color_hex)

            # LAYER 2: Vector Graphics Reconstruction (Round corner shapes, Cards, Lines)
            shapes_list = slide_def.get("shapes", [])
            for shape_def in shapes_list:
                shape_type_str = shape_def.get("type", "rectangle")
                
                # Match to python-pptx built-in enum types
                if shape_type_str == "round-rect":
                    mso_type = MSO_SHAPE.ROUNDED_RECTANGLE
                elif shape_type_str == "ellipse":
                    mso_type = MSO_SHAPE.OVAL
                elif shape_type_str == "line":
                    mso_type = MSO_SHAPE.RECTANGLE # In pptx, lines are best represented as filled rectangles for thickness stability
                else: 
                    mso_type = MSO_SHAPE.RECTANGLE

                # Coordinate adjustments utilizing calibration offsets
                x = self.offsets.adjust_x(shape_def.get("x", 0))
                y = self.offsets.adjust_y(shape_def.get("y", 0))
                w = self.offsets.adjust_width(shape_def.get("w", 10))
                h = self.offsets.adjust_height(shape_def.get("h", 10))

                shape = slide.shapes.add_shape(mso_type, x, y, w, h)
                
                # Style filling
                fill_color_hex = shape_def.get("fill", "E5E7EB")
                shape.fill.solid()
                shape.fill.fore_color.rgb = self._hex_to_rgb(fill_color_hex)

                # Line border styling
                border_def = shape_def.get("border")
                if border_def:
                    shape.line.color.rgb = self._hex_to_rgb(border_def.get("color", "BDC3C7"))
                    shape.line.width = Pt(border_def.get("width", 1.0))
                else:
                    shape.line.fill.background() # No line border

            # LAYER 3: Typography Text Layout Rebuilding
            texts_list = slide_def.get("texts", [])
            for text_def in texts_list:
                tx_x = self.offsets.adjust_x(text_def.get("x", 0))
                tx_y = self.offsets.adjust_y(text_def.get("y", 0))
                tx_w = self.offsets.adjust_width(text_def.get("w", 20))
                tx_h = self.offsets.adjust_height(text_def.get("h", 10))

                textbox = slide.shapes.add_textbox(tx_x, tx_y, tx_w, tx_h)
                tf = textbox.text_frame
                tf.word_wrap = True
                
                p = tf.paragraphs[0]
                p.text = text_def.get("text", "")
                
                # Apply Font configurations
                p.font.name = "Arial"
                p.font.size = Pt(text_def.get("fontSize", 12))
                p.font.color.rgb = self._hex_to_rgb(text_def.get("color", "1F2937"))
                p.font.bold = bool(text_def.get("bold", False))
                p.font.italic = bool(text_def.get("italic", False))

                # Alignments matching standard formats
                align_str = text_def.get("alignment", "left").lower()
                if align_str == "center":
                    p.alignment = PP_ALIGN.CENTER
                elif align_str == "right":
                    p.alignment = PP_ALIGN.RIGHT
                else:
                    p.alignment = PP_ALIGN.LEFT

        return prs
